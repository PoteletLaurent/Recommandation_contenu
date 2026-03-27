"""
Génère le support de présentation PowerPoint du projet P10.
Usage : python3 generate_presentation.py
Output : presentation_p10.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── Palette ──────────────────────────────────────────────
BLEU_FONCE  = RGBColor(0x1F, 0x3A, 0x5F)   # titres
BLEU_CLAIR  = RGBColor(0x2E, 0x86, 0xC1)   # accents
GRIS_FOND   = RGBColor(0xF4, 0xF6, 0xF9)
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
VERT        = RGBColor(0x1A, 0x8A, 0x5A)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
ROUGE       = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout vide


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def text_box(slide, text, x, y, w, h, size=18, bold=False, color=BLEU_FONCE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def bullet_box(slide, items, x, y, w, h, size=16, color=BLEU_FONCE, indent=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("      " if indent else "  ") + item
        p.font.size = Pt(size)
        p.font.color.rgb = color


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.2, fill=BLEU_FONCE)
    text_box(slide, title, 0.4, 0.15, 12, 0.75, size=28, bold=True,
             color=BLANC, align=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, subtitle, 0.4, 0.82, 12, 0.4, size=14,
                 color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.LEFT)


def footer(slide, num):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=BLEU_FONCE)
    text_box(slide, f"Système de recommandation d'articles  |  MVP", 0.3, 7.22,
             11, 0.25, size=10, color=BLANC)
    text_box(slide, str(num), 12.8, 7.22, 0.4, 0.25, size=10,
             color=BLANC, align=PP_ALIGN.RIGHT)


def accent_bar(slide):
    rect(slide, 0, 1.2, 0.08, 6.0, fill=BLEU_CLAIR)


# ══════════════════════════════════════════════════════════
# SLIDE 1 – Couverture
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=BLEU_FONCE)
rect(sl, 0, 0, 13.33, 0.15, fill=BLEU_CLAIR)
rect(sl, 0, 7.35, 13.33, 0.15, fill=BLEU_CLAIR)

text_box(sl, "Système de recommandation d'articles", 1, 1.8, 11.33, 1.2,
         size=36, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
text_box(sl, "MVP – Architecture serverless sur AWS", 1, 3.0, 11.33, 0.6,
         size=22, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)
text_box(sl, "Données : News Portal User Interactions – Globo.com", 1, 3.8,
         11.33, 0.5, size=16, color=RGBColor(0x85, 0xC1, 0xE9), align=PP_ALIGN.CENTER)
text_box(sl, "2026", 1, 6.5, 11.33, 0.5, size=14,
         color=RGBColor(0x7F, 0x8C, 0x8D), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2 – Sommaire
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Sommaire")
accent_bar(sl)
footer(sl, 2)

sections = [
    ("1", "Contexte et objectif"),
    ("2", "Données utilisées"),
    ("3", "Description fonctionnelle de l'application"),
    ("4", "Architecture technique retenue"),
    ("5", "Système de recommandation – 3 modèles"),
    ("6", "Comparaison et choix des modèles"),
    ("7", "Architecture cible – Gestion du cold start"),
    ("8", "Limites et perspectives"),
]

for i, (num, label) in enumerate(sections):
    row = i % 4
    col = i // 4
    x = 1.0 + col * 6.2
    y = 1.5 + row * 1.3
    rect(sl, x, y, 5.8, 1.0, fill=BLEU_FONCE)
    text_box(sl, num, x + 0.15, y + 0.1, 0.6, 0.8, size=28, bold=True,
             color=BLEU_CLAIR, align=PP_ALIGN.CENTER)
    text_box(sl, label, x + 0.8, y + 0.2, 4.8, 0.7, size=15,
             color=BLANC, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════
# SLIDE 3 – Contexte et objectif
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Contexte et objectif", "Pourquoi ce projet ?")
accent_bar(sl)
footer(sl, 3)

rect(sl, 0.5, 1.4, 5.8, 5.5, fill=BLANC)
text_box(sl, "Contexte", 0.7, 1.5, 5.4, 0.5, size=16, bold=True, color=BLEU_CLAIR)
bullet_box(sl, [
    "Tester une solution de recommandation d'articles pour des particuliers",
    "Pas encore de données utilisateurs propriétaires",
    "Utilisation d'un dataset public pour développer le MVP",
], 0.7, 1.95, 5.4, 2.5, size=14)

rect(sl, 7.0, 1.4, 5.8, 5.5, fill=BLANC)
text_box(sl, "Fonctionnalité MVP", 7.2, 1.5, 5.4, 0.5, size=16, bold=True, color=BLEU_CLAIR)
text_box(sl, "« En tant qu'utilisateur, je reçois une sélection de 5 articles recommandés »",
         7.2, 2.0, 5.4, 1.2, size=14, color=BLEU_FONCE)
text_box(sl, "Enjeu architectural clé", 7.2, 3.4, 5.4, 0.4, size=14, bold=True, color=BLEU_FONCE)
bullet_box(sl, [
    "Prise en compte des nouveaux utilisateurs",
    "Prise en compte des nouveaux articles",
    "Architecture scalable et maintenable",
], 7.2, 3.8, 5.4, 2.0, size=14)

# ══════════════════════════════════════════════════════════
# SLIDE 4 – Données
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Données utilisées", "News Portal User Interactions – Globo.com (Kaggle, 2017)")
accent_bar(sl)
footer(sl, 4)

fichiers = [
    ("articles_metadata.csv", BLEU_CLAIR,
     ["364 047 articles", "Colonnes : article_id, category_id, words_count",
      "Dataset anonymisé (pas de titres)", "461 catégories distinctes"]),
    ("clicks_sample.csv", VERT,
     ["1 883 interactions", "707 utilisateurs uniques", "323 articles uniques cliqués",
      "Colonnes : user_id, session, device, pays..."]),
    ("articles_embeddings.pickle", ORANGE,
     ["Vecteurs de 250 dimensions par article", "Générés depuis le contenu textuel",
      "Fichier brut : 347 Mo", "Après PCA (64D) : 93 Mo"]),
]

for i, (nom, couleur, details) in enumerate(fichiers):
    x = 0.4 + i * 4.3
    rect(sl, x, 1.4, 4.1, 0.55, fill=couleur)
    text_box(sl, nom, x + 0.15, 1.48, 3.8, 0.4, size=13, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    rect(sl, x, 1.95, 4.1, 4.8, fill=BLANC)
    bullet_box(sl, details, x + 0.2, 2.1, 3.7, 4.4, size=13)

# ══════════════════════════════════════════════════════════
# SLIDE 5 – Description fonctionnelle
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Description fonctionnelle de l'application")
accent_bar(sl)
footer(sl, 5)

text_box(sl, "L'interface Streamlit permet à l'utilisateur de :", 0.5, 1.4, 12.5, 0.5,
         size=16, bold=True, color=BLEU_FONCE)

etapes = [
    ("1", "Sélectionner un utilisateur", "Choisir parmi les utilisateurs connus via une liste déroulante"),
    ("2", "Choisir un modèle", "ALS, Embeddings PCA ou Similarité item-based"),
    ("3", "Lancer la recommandation", "Appel de la fonction AWS Lambda via boto3"),
    ("4", "Visualiser les résultats", "5 articles avec catégorie, longueur et indicateur de cohérence"),
]

for i, (num, titre, desc) in enumerate(etapes):
    x = 0.4 + (i % 2) * 6.4
    y = 2.0 + (i // 2) * 2.3
    rect(sl, x, y, 6.0, 2.0, fill=BLANC)
    rect(sl, x, y, 0.6, 2.0, fill=BLEU_CLAIR)
    text_box(sl, num, x + 0.05, y + 0.55, 0.5, 0.8, size=24, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    text_box(sl, titre, x + 0.75, y + 0.2, 5.0, 0.5, size=15, bold=True, color=BLEU_FONCE)
    text_box(sl, desc, x + 0.75, y + 0.75, 5.0, 1.0, size=13, color=BLEU_FONCE)

# ══════════════════════════════════════════════════════════
# SLIDE 6 – Architecture retenue schéma
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Architecture technique retenue", "Serverless sur AWS – sans API Gateway")
accent_bar(sl)
footer(sl, 6)

composants = [
    (0.5,  3.2, 2.2, 1.2, BLEU_FONCE, "Utilisateur", "Navigateur web"),
    (3.5,  3.2, 2.5, 1.2, BLEU_CLAIR, "Streamlit", "EC2 – port 8501"),
    (7.0,  3.2, 2.5, 1.2, VERT,       "AWS Lambda", "p10-recommandation-handler"),
    (10.1, 3.2, 2.6, 1.2, ORANGE,     "AWS S3", "p10-recommandation-models"),
]

for (x, y, w, h, col, titre, sous) in composants:
    rect(sl, x, y, w, h, fill=col)
    text_box(sl, titre, x + 0.1, y + 0.1, w - 0.2, 0.5, size=15, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    text_box(sl, sous, x + 0.1, y + 0.6, w - 0.2, 0.4, size=11,
             color=BLANC, align=PP_ALIGN.CENTER)

fleches = [
    (2.7, 3.78, "HTTP"),
    (6.0, 3.78, "boto3"),
    (9.5, 3.78, "S3 download"),
]
for (x, y, label) in fleches:
    rect(sl, x, y - 0.02, 0.8, 0.08, fill=BLEU_FONCE)
    text_box(sl, "▶", x + 0.65, y - 0.18, 0.3, 0.3, size=14, color=BLEU_FONCE)
    text_box(sl, label, x, y + 0.1, 0.8, 0.3, size=10,
             color=BLEU_FONCE, align=PP_ALIGN.CENTER)

text_box(sl, "Artefacts stockés sur S3 :", 0.5, 5.2, 4.5, 0.4,
         size=13, bold=True, color=BLEU_FONCE)
bullet_box(sl, ["mappings.pkl", "als_model.pkl", "embeddings_pca.pkl",
                "item_similarity.pkl", "user_clicks.pkl"],
           0.5, 5.55, 4.0, 1.5, size=12)

text_box(sl, "Choix de conception :", 5.5, 5.2, 7.3, 0.4,
         size=13, bold=True, color=BLEU_FONCE)
bullet_box(sl, [
    "Pas d'API Gateway : appel direct boto3 (plus simple pour un MVP)",
    "Lambda mise en cache les artefacts en mémoire (cold start unique)",
    "3 modèles sélectionnables dynamiquement",
], 5.5, 5.55, 7.3, 1.5, size=12)

# ══════════════════════════════════════════════════════════
# SLIDE 7 – Modèle 1 : ALS
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Modèle 1 – ALS (Alternating Least Squares)", "Filtrage collaboratif")
accent_bar(sl)
footer(sl, 7)

rect(sl, 0.5, 1.4, 8.0, 5.5, fill=BLANC)
text_box(sl, "Principe", 0.7, 1.5, 7.6, 0.4, size=15, bold=True, color=BLEU_CLAIR)
text_box(sl,
    "Factorise la matrice utilisateurs × articles en deux matrices de facteurs latents. "
    "Les recommandations sont les articles dont les facteurs sont les plus proches de ceux de l'utilisateur.",
    0.7, 1.9, 7.6, 1.0, size=13, color=BLEU_FONCE)

text_box(sl, "Avantages", 0.7, 3.1, 3.5, 0.4, size=14, bold=True, color=VERT)
bullet_box(sl, ["Capture les préférences implicites", "Efficace sur grandes matrices creuses",
                "Bien documenté et éprouvé"], 0.7, 3.5, 3.5, 2.0, size=13)

text_box(sl, "Inconvénients", 4.5, 3.1, 3.8, 0.4, size=14, bold=True, color=ROUGE)
bullet_box(sl, ["Cold start : échoue sans historique", "Nécessite un réentraînement régulier",
                "Ne tient pas compte du contenu"], 4.5, 3.5, 3.8, 2.0, size=13)

rect(sl, 9.0, 1.4, 3.8, 5.5, fill=BLEU_FONCE)
text_box(sl, "Résultats", 9.2, 1.5, 3.4, 0.4, size=15, bold=True, color=BLANC)
for label, val, y in [("Rappel@10", "0.11", 2.2), ("Précision@10", "0.02", 3.3), ("F1-Score", "0.03", 4.4)]:
    text_box(sl, label, 9.2, y, 3.4, 0.4, size=13, color=RGBColor(0xAD, 0xD8, 0xE6),
             align=PP_ALIGN.CENTER)
    text_box(sl, val, 9.2, y + 0.4, 3.4, 0.6, size=28, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 8 – Modèle 2 : Embeddings PCA
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Modèle 2 – Embeddings + PCA", "Recommandation basée sur le contenu")
accent_bar(sl)
footer(sl, 8)

rect(sl, 0.5, 1.4, 8.0, 5.5, fill=BLANC)
text_box(sl, "Principe", 0.7, 1.5, 7.6, 0.4, size=15, bold=True, color=BLEU_CLAIR)
text_box(sl,
    "Chaque article est représenté par un vecteur de 250 dimensions. "
    "Une PCA réduit ces vecteurs à 64D (97.21% de variance conservée). "
    "Le profil utilisateur est la moyenne des vecteurs de ses articles lus. "
    "On recommande les articles les plus proches par similarité cosinus.",
    0.7, 1.9, 7.6, 1.3, size=13, color=BLEU_FONCE)

text_box(sl, "Avantages", 0.7, 3.4, 3.5, 0.4, size=14, bold=True, color=VERT)
bullet_box(sl, ["Fonctionne pour les nouveaux articles", "Ne nécessite pas d'historique de clics",
                "Capture la sémantique du contenu"], 0.7, 3.8, 3.5, 2.0, size=13)

text_box(sl, "Inconvénients", 4.5, 3.4, 3.8, 0.4, size=14, bold=True, color=ROUGE)
bullet_box(sl, ["Cold start utilisateur toujours présent", "Qualité dépend des embeddings source",
                "Fichier lourd (347Mo → 93Mo après PCA)"], 4.5, 3.8, 3.8, 2.0, size=13)

rect(sl, 9.0, 1.4, 3.8, 5.5, fill=BLEU_FONCE)
text_box(sl, "Résultats", 9.2, 1.5, 3.4, 0.4, size=15, bold=True, color=BLANC)
for label, val, y in [("Rappel@10", "0.21", 2.2), ("Précision@10", "0.03", 3.3), ("F1-Score", "0.06", 4.4)]:
    text_box(sl, label, 9.2, y, 3.4, 0.4, size=13, color=RGBColor(0xAD, 0xD8, 0xE6),
             align=PP_ALIGN.CENTER)
    text_box(sl, val, 9.2, y + 0.4, 3.4, 0.6, size=28, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)

text_box(sl, "PCA : 250D → 64D\n97.21% variance conservée", 9.2, 5.6, 3.4, 0.8,
         size=12, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 9 – Modèle 3 : Similarité item-based
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Modèle 3 – Similarité item-based", "Filtrage collaboratif basé sur les articles")
accent_bar(sl)
footer(sl, 9)

rect(sl, 0.5, 1.4, 8.0, 5.5, fill=BLANC)
text_box(sl, "Principe", 0.7, 1.5, 7.6, 0.4, size=15, bold=True, color=BLEU_CLAIR)
text_box(sl,
    "Calcule la similarité cosinus entre tous les articles à partir des interactions utilisateurs. "
    "Pour un utilisateur, on additionne les scores de similarité de ses articles lus "
    "et on recommande les articles au score cumulé le plus élevé.",
    0.7, 1.9, 7.6, 1.2, size=13, color=BLEU_FONCE)

text_box(sl, "Avantages", 0.7, 3.3, 3.5, 0.4, size=14, bold=True, color=VERT)
bullet_box(sl, ["Meilleurs résultats sur ce dataset", "Simple à interpréter",
                "Pas de réentraînement lourd requis"], 0.7, 3.7, 3.5, 2.0, size=13)

text_box(sl, "Inconvénients", 4.5, 3.3, 3.8, 0.4, size=14, bold=True, color=ROUGE)
bullet_box(sl, ["Cold start utilisateur et article", "Matrice de similarité à recalculer",
                "Limité aux articles déjà vus"], 4.5, 3.7, 3.8, 2.0, size=13)

rect(sl, 9.0, 1.4, 3.8, 5.5, fill=VERT)
text_box(sl, "Résultats", 9.2, 1.5, 3.4, 0.4, size=15, bold=True, color=BLANC)
for label, val, y in [("Rappel@10", "0.34", 2.2), ("Précision@10", "0.04", 3.3), ("F1-Score", "0.07", 4.4)]:
    text_box(sl, label, 9.2, y, 3.4, 0.4, size=13, color=RGBColor(0xD5, 0xF5, 0xE3),
             align=PP_ALIGN.CENTER)
    text_box(sl, val, 9.2, y + 0.4, 3.4, 0.6, size=28, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)

text_box(sl, "MEILLEUR MODÈLE", 9.2, 5.6, 3.4, 0.5, size=13, bold=True,
         color=BLANC, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 10 – Comparaison des modèles
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Comparaison des modèles")
accent_bar(sl)
footer(sl, 10)

cols_labels = ["Modèle", "Type", "Rappel@10", "F1-Score", "Cold start\nutilisateur", "Cold start\narticle"]
rows = [
    ["ALS",                   "Collaborative",  "0.11", "0.03", "✗", "✗"],
    ["Embeddings PCA",        "Content-based",  "0.21", "0.06", "✗", "✓"],
    ["Similarité item-based", "Item-based",     "0.34", "0.07", "✗", "✗"],
]
col_widths = [2.4, 2.2, 1.6, 1.5, 1.9, 1.9]
col_colors_header = [BLEU_FONCE] * 6
row_fills = [BLANC, BLANC, VERT]
text_colors_rows = [BLEU_FONCE, BLEU_FONCE, BLANC]

x_start, y_start = 0.5, 1.5
row_h = 0.85

x = x_start
for j, (label, w) in enumerate(zip(cols_labels, col_widths)):
    rect(sl, x, y_start, w - 0.05, row_h * 0.9, fill=BLEU_FONCE)
    text_box(sl, label, x + 0.05, y_start + 0.08, w - 0.15, row_h * 0.8,
             size=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    x += w

for i, (row, fill, tcol) in enumerate(zip(rows, row_fills, text_colors_rows)):
    y = y_start + (i + 1) * row_h
    x = x_start
    for j, (val, w) in enumerate(zip(row, col_widths)):
        rect(sl, x, y, w - 0.05, row_h * 0.9, fill=fill)
        vc = VERT if (val == "✓") else (ROUGE if val == "✗" else tcol)
        text_box(sl, val, x + 0.05, y + 0.15, w - 0.15, row_h * 0.65,
                 size=14, bold=(j == 0), color=vc, align=PP_ALIGN.CENTER)
        x += w

text_box(sl,
    "Le modèle de similarité item-based offre les meilleures performances sur ce dataset. "
    "Les 3 modèles sont accessibles dynamiquement depuis l'interface.",
    0.5, 6.0, 12.5, 0.8, size=14, color=BLEU_FONCE)

# ══════════════════════════════════════════════════════════
# SLIDE 11 – Architecture cible : problématique
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Architecture cible – Problématique du cold start")
accent_bar(sl)
footer(sl, 11)

text_box(sl, "Le cold start désigne l'incapacité à recommander sans historique.",
         0.5, 1.4, 12.5, 0.5, size=15, bold=True, color=BLEU_FONCE)

cas = [
    (ROUGE, "Nouvel utilisateur",
     ["Aucun clic enregistré", "ALS et similarité item-based échouent",
      "Solution : popularité globale + onboarding (préférences déclarées)"]),
    (ORANGE, "Nouvel article",
     ["Jamais consulté par personne", "ALS et similarité item-based échouent",
      "Solution : embeddings générés depuis le contenu à la publication"]),
    (BLEU_CLAIR, "Utilisateur avec peu d'historique",
     ["Moins de 3 articles lus", "Recommandations peu fiables",
      "Solution : approche hybride (embeddings + popularité)"]),
]

for i, (col, titre, details) in enumerate(cas):
    x = 0.4 + i * 4.3
    rect(sl, x, 2.1, 4.1, 0.55, fill=col)
    text_box(sl, titre, x + 0.1, 2.18, 3.9, 0.4, size=14, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    rect(sl, x, 2.65, 4.1, 4.0, fill=BLANC)
    bullet_box(sl, details, x + 0.15, 2.8, 3.8, 3.5, size=13)

# ══════════════════════════════════════════════════════════
# SLIDE 12 – Architecture cible : schéma
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Architecture cible – Schéma")
accent_bar(sl)
footer(sl, 12)

rect(sl, 0.3, 1.3, 12.7, 5.6, fill=BLANC)

text_box(sl, "Streamlit", 0.5, 1.5, 2.5, 0.5, size=13, bold=True,
         color=BLEU_CLAIR, align=PP_ALIGN.CENTER)
rect(sl, 0.5, 2.0, 2.5, 0.7, fill=BLEU_CLAIR)
text_box(sl, "user_id + contexte", 0.55, 2.1, 2.4, 0.5, size=11,
         color=BLANC, align=PP_ALIGN.CENTER)

text_box(sl, "Lambda – Routeur", 3.5, 1.5, 2.8, 0.5, size=13, bold=True,
         color=BLEU_FONCE, align=PP_ALIGN.CENTER)
rect(sl, 3.5, 2.0, 2.8, 0.7, fill=BLEU_FONCE)
text_box(sl, "Détecte le cas", 3.55, 2.1, 2.7, 0.5, size=11,
         color=BLANC, align=PP_ALIGN.CENTER)

stratégies = [
    (0.5,  3.5, VERT,       "User connu",      "ALS /\nSimilarité"),
    (3.3,  3.5, BLEU_CLAIR, "Peu d'historique","Hybride\nEmb. + Pop."),
    (6.1,  3.5, ORANGE,     "Nouvel user",      "Onboarding\n+ Popularité"),
    (8.9,  3.5, ROUGE,      "Nouvel article",   "Embeddings\ncontenu"),
]
for (x, y, col, label, action) in stratégies:
    rect(sl, x, y, 2.5, 0.5, fill=col)
    text_box(sl, label, x + 0.05, y + 0.08, 2.4, 0.35, size=11, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    rect(sl, x, y + 0.5, 2.5, 0.9, fill=BLANC)
    text_box(sl, action, x + 0.05, y + 0.58, 2.4, 0.75, size=12,
             color=BLEU_FONCE, align=PP_ALIGN.CENTER)

text_box(sl, "EventBridge (24h)", 0.5, 5.3, 3.0, 0.4, size=12, bold=True, color=BLEU_FONCE)
text_box(sl, "→ Réentraînement automatique → Mise à jour S3 → Lambda recharge les artefacts",
         0.5, 5.7, 12.0, 0.5, size=12, color=BLEU_FONCE)

# ══════════════════════════════════════════════════════════
# SLIDE 13 – Réentraînement continu
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Réentraînement continu", "Intégrer les nouveaux utilisateurs et articles au fil du temps")
accent_bar(sl)
footer(sl, 13)

etapes = [
    (BLEU_CLAIR, "1. Nouveau contenu",       "Nouvel article publié → génération automatique de son embedding via un modèle NLP léger → ajout dans S3"),
    (VERT,       "2. Accumulation des clics", "Les nouvelles interactions utilisateurs s'accumulent dans la base de clics"),
    (ORANGE,     "3. Déclenchement",          "AWS EventBridge déclenche un job de réentraînement toutes les 24h"),
    (BLEU_FONCE, "4. Réentraînement",         "Les 3 modèles sont réentraînés sur les données mises à jour, les nouveaux artefacts sont uploadés sur S3"),
    (VERT,       "5. Mise en production",     "La Lambda recharge automatiquement les artefacts au prochain appel (cache invalidé)"),
]

for i, (col, titre, desc) in enumerate(etapes):
    y = 1.5 + i * 1.0
    rect(sl, 0.5, y, 0.7, 0.75, fill=col)
    text_box(sl, str(i+1), 0.5, y + 0.12, 0.7, 0.5, size=18, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    rect(sl, 1.3, y, 11.5, 0.75, fill=BLANC)
    text_box(sl, titre, 1.5, y + 0.05, 3.0, 0.4, size=13, bold=True, color=col)
    text_box(sl, desc, 4.7, y + 0.1, 7.9, 0.55, size=12, color=BLEU_FONCE)

# ══════════════════════════════════════════════════════════
# SLIDE 14 – Limites et perspectives
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_FOND)
header(sl, "Limites du MVP et perspectives")
accent_bar(sl)
footer(sl, 14)

limites = [
    ("Dataset anonymisé",          "Pas de titres d'articles → affichage peu parlant",     "Remplacer par des données propriétaires avec titres"),
    ("Cold start non géré",        "Nouveaux users/articles sans recommandation",           "Implémenter les stratégies de l'architecture cible"),
    ("Pas de réentraînement auto", "Modèle statique après déploiement",                     "EventBridge + Lambda de réentraînement"),
    ("Streamlit non persistant",   "Redémarrage EC2 = perte du service",                   "Configurer systemd ou supervisor pour relance auto"),
]

rect(sl, 0.3, 1.4, 3.8, 0.5, fill=ROUGE)
rect(sl, 4.5, 1.4, 4.3, 0.5, fill=ORANGE)
rect(sl, 9.2, 1.4, 3.7, 0.5, fill=VERT)
for x, label, col in [(0.3, "Limite", ROUGE), (4.5, "Impact", ORANGE), (9.2, "Solution cible", VERT)]:
    text_box(sl, label, x + 0.1, 1.48, 3.5, 0.35, size=13, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)

for i, (limite, impact, solution) in enumerate(limites):
    y = 2.1 + i * 1.15
    fill = GRIS_FOND if i % 2 == 0 else BLANC
    for x, w, txt in [(0.3, 3.8, limite), (4.5, 4.3, impact), (9.2, 3.7, solution)]:
        rect(sl, x, y, w, 1.0, fill=fill)
        text_box(sl, txt, x + 0.15, y + 0.15, w - 0.25, 0.75, size=12, color=BLEU_FONCE)

# ══════════════════════════════════════════════════════════
# SLIDE 15 – Conclusion
# ══════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=BLEU_FONCE)
rect(sl, 0, 0, 13.33, 0.12, fill=BLEU_CLAIR)
rect(sl, 0, 7.38, 13.33, 0.12, fill=BLEU_CLAIR)

text_box(sl, "Conclusion", 1, 0.8, 11.33, 0.8, size=32, bold=True,
         color=BLANC, align=PP_ALIGN.CENTER)

points = [
    ("MVP fonctionnel",         "Application Streamlit déployée sur EC2, 3 modèles accessibles, Lambda opérationnelle"),
    ("Architecture serverless", "AWS Lambda + S3 : scalable, économique, sans serveur à maintenir"),
    ("Meilleur modèle",         "Similarité item-based : Rappel@10 = 0.34, meilleure performance sur ce dataset"),
    ("Industrialisable",        "Scripts versionnés sur GitHub, déploiement reproductible bout-en-bout"),
    ("Architecture cible",      "Cold start géré par routage intelligent selon le contexte utilisateur/article"),
]

for i, (titre, desc) in enumerate(points):
    y = 1.7 + i * 1.0
    rect(sl, 1.0, y, 0.5, 0.7, fill=BLEU_CLAIR)
    text_box(sl, "✓", 1.0, y + 0.1, 0.5, 0.5, size=16, bold=True,
             color=BLANC, align=PP_ALIGN.CENTER)
    text_box(sl, titre, 1.7, y + 0.05, 2.8, 0.4, size=14, bold=True, color=BLEU_CLAIR)
    text_box(sl, desc, 4.7, y + 0.1, 8.0, 0.5, size=13, color=RGBColor(0xD6, 0xEA, 0xF8))

# ══════════════════════════════════════════════════════════
# Sauvegarde
# ══════════════════════════════════════════════════════════
output_path = "docs/presentation_p10.pptx"
prs.save(output_path)
print(f"Présentation générée : {output_path}")
print(f"Nombre de slides : {len(prs.slides)}")
